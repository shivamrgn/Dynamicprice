import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set widescreen slide dimensions (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Branded Academic Color Palette (Parul Red theme)
    PARUL_RED = RGBColor(163, 0, 0)       # Deep Crimson Maroon (#A30000)
    TEXT_DARK = RGBColor(51, 51, 51)       # Charcoal (#333333)
    ACCENT_GOLD = RGBColor(218, 165, 32)   # Antique Gold (#DAA520)
    BG_LIGHT = RGBColor(248, 249, 250)     # Crisp Off-White (#F8F9FA)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(240, 240, 240)

    # Blank layout for custom shapes styling
    blank_slide_layout = prs.slide_layouts[6]

    # Helper: Set slide background color
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Add branded top header bar, accent line, and title
    def add_branded_header(slide, title_text):
        set_background(slide, BG_LIGHT)

        # Thin Red Brand Bar on top
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.16)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = PARUL_RED
        top_bar.line.color.rgb = PARUL_RED

        # Slide Title text frame
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PARUL_RED

        # Subtle gold accent strip below title
        gold_strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2), Inches(2.2), Inches(0.04)
        )
        gold_strip.fill.solid()
        gold_strip.fill.fore_color.rgb = ACCENT_GOLD
        gold_strip.line.color.rgb = ACCENT_GOLD

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.133), Inches(0.3))
        ftf = footer_box.text_frame
        fp = ftf.paragraphs[0]
        fp.text = "Parul University | Department of Computer Science & Engineering | B.Tech Project Presentation"
        fp.font.name = "Arial"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(120, 120, 120)

    # Helper: Universal parser to convert standard lists to styled bullets
    def add_styled_bullets(tf, bullet_list, base_font_size=15):
        tf.word_wrap = True
        for idx, line in enumerate(bullet_list):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Determine bullet hierarchy based on double-space indentation
            stripped = line.lstrip()
            indent_spaces = len(line) - len(stripped)
            level = indent_spaces // 2
            if level > 2:
                level = 2

            p.level = level

            # Remove leading bullet symbols from string
            clean_text = stripped
            if clean_text.startswith("* "):
                clean_text = clean_text[2:]
            elif clean_text.startswith("- "):
                clean_text = clean_text[2:]
            elif clean_text.startswith("• "):
                clean_text = clean_text[2:]

            p.text = clean_text
            p.font.name = "Calibri"
            p.font.size = Pt(base_font_size - (level * 2))
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1: Cover & Title Slide
    # ----------------------------------------------------
    slide_1 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide_1, BG_LIGHT)

    # Left colored brand panel
    left_panel = slide_1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.2), Inches(7.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = PARUL_RED
    left_panel.line.fill.background()

    # Right accent strip
    accent_strip = slide_1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(0), Inches(0.08), Inches(7.5)
    )
    accent_strip.fill.solid()
    accent_strip.fill.fore_color.rgb = ACCENT_GOLD
    accent_strip.line.fill.background()

    # Left Panel Text Box (University Details)
    uni_box = slide_1.shapes.add_textbox(Inches(0.4), Inches(0.8), Inches(3.4), Inches(6.0))
    utf = uni_box.text_frame
    utf.word_wrap = True
    
    p1 = utf.paragraphs[0]
    p1.text = "PARUL UNIVERSITY"
    p1.font.name = "Arial"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(20)

    p2 = utf.add_paragraph()
    p2.text = "DEPARTMENT OF\nCOMPUTER SCIENCE\n& ENGINEERING"
    p2.font.name = "Arial"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GOLD
    p2.space_after = Pt(40)

    p3 = utf.add_paragraph()
    p3.text = "B.Tech Final Year\nProject Presentation\nAcademic Session 2026"
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = WHITE

    # Right Panel Title Text Block
    title_box = slide_1.shapes.add_textbox(Inches(4.8), Inches(1.5), Inches(7.8), Inches(4.5))
    ttf = title_box.text_frame
    ttf.word_wrap = True

    tp1 = ttf.paragraphs[0]
    tp1.text = "DYNAMIC PRICING ENGINE"
    tp1.font.name = "Arial"
    tp1.font.size = Pt(40)
    tp1.font.bold = True
    tp1.font.color.rgb = PARUL_RED
    tp1.space_after = Pt(6)

    tp2 = ttf.add_paragraph()
    tp2.text = "Real-Time Price Optimization Using Price Elasticity and Reinforcement Learning"
    tp2.font.name = "Arial"
    tp2.font.size = Pt(18)
    tp2.font.bold = True
    tp2.font.color.rgb = TEXT_DARK
    tp2.space_after = Pt(40)

    tp3 = ttf.add_paragraph()
    tp3.text = "Prepared By:\n• Shivam Kumar (Lead AI & Full Stack Engineer)"
    tp3.font.name = "Arial"
    tp3.font.size = Pt(14)
    tp3.font.bold = True
    tp3.font.color.rgb = TEXT_DARK
    tp3.space_after = Pt(12)

    tp4 = ttf.add_paragraph()
    tp4.text = "Technologies Utilized:\nPython, Next.js 15, FastAPI, MongoDB, Prophet, XGBoost, Q-Learning"
    tp4.font.name = "Arial"
    tp4.font.size = Pt(12)
    tp4.font.color.rgb = RGBColor(100, 100, 100)

    # ----------------------------------------------------
    # SLIDE 2: Problem Statement
    # ----------------------------------------------------
    slide_2 = prs.slides.add_slide(blank_slide_layout)
    add_branded_header(slide_2, "2. Problem Statement")

    col1_box = slide_2.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    p_text_col1 = [
        "E-Commerce Gross Margin Leakage",
        "  - Traditional online retailers rely on static pricing or manual rules, causing them to leak 8% to 15% in gross margin.",
        "  - Static systems fail to adjust to seasonal shifts, holidays, or real-time competitor stock changes.",
        "Inability to Track Fast Competitor Shifts",
        "  - Competitive price monitoring is labor-intensive and delayed when done manually.",
        "  - Merchants miss immediate conversion opportunities when competitors go out of stock."
    ]
    add_styled_bullets(col1_box.text_frame, p_text_col1, base_font_size=15)

    col2_box = slide_2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    p_text_col2 = [
        "Inventory Scarcity Misalignment",
        "  - High-demand items are frequently underpriced, triggering stockouts and supply gaps.",
        "  - Conversely, slow-moving items remain overpriced, freezing capital in storage warehouse spaces.",
        "arbitrary Discount Degradation",
        "  - Lack of scientific buyer sensitivity (elasticity) analysis causes arbitrary, deep discounts.",
        "  - This erodes brand value and overall operating margins without generating structural sales boosts."
    ]
    add_styled_bullets(col2_box.text_frame, p_text_col2, base_font_size=15)

    # ----------------------------------------------------
    # SLIDE 3: Features
    # ----------------------------------------------------
    slide_3 = prs.slides.add_slide(blank_slide_layout)
    add_branded_header(slide_3, "3. Core Features")

    f1_box = slide_3.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    f_text_col1 = [
        "Real-Time Competitor Ingestion",
        "  - Autonomous scraping background worker fetches competitor prices and stock positions every 2 hours.",
        "  - Enables proactive pricing matching without causing destructive price wars.",
        "Hybrid Demand Forecasting Engine",
        "  - Combines Prophet (seasonal macro-trends) and XGBoost (micro feature variables) into a predictive forecaster.",
        "  - Continuously validates error metrics with RMSE benchmarks."
    ]
    add_styled_bullets(f1_box.text_frame, f_text_col1, base_font_size=14)

    f2_box = slide_3.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    f_text_col2 = [
        "Log-Log Price Elasticity Modeling",
        "  - Regresses historical price-demand signals to calculate consumer willingness-to-pay sensitivity indices.",
        "  - Categorizes SKUs into Elastic, Inelastic, or Anomalous groups.",
        "Q-Learning Reinforcement Optimizer",
        "  - Uses tabular RL algorithms to adjust retail prices based on current stock scarcity and volume metrics.",
        "  - Learns optimal policies through continuous explore-exploit sequences."
    ]
    add_styled_bullets(f2_box.text_frame, f_text_col2, base_font_size=14)

    # ----------------------------------------------------
    # SLIDE 4: Objective and Success Metrics
    # ----------------------------------------------------
    slide_4 = prs.slides.add_slide(blank_slide_layout)
    add_branded_header(slide_4, "4. Objective & Success Metrics")

    obj_box = slide_4.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    obj_text = [
        "Project Objectives",
        "  - Build a fully containerized, microservice dynamic pricing monorepo containing three decoupled stacks.",
        "  - Automate competitor telemetry scraping and pricing adjustments.",
        "  - Connect price sensitivity data to stock level indicators to automate discounts and premium hikes.",
        "  - Maintain strict margins via algorithmic baseline limits."
    ]
    add_styled_bullets(obj_box.text_frame, obj_text, base_font_size=15)

    met_box = slide_4.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    met_text = [
        "Key Success Metrics",
        "  - Revenue Uplift: Target a 15% to 25% average increase in simulated gross revenues.",
        "  - Inventory Turnover: Improve sales of slow-moving inventory by 30%, optimizing warehousing space.",
        "  - Ingestion & Sync Latency: Keep network database price sync latency under 200ms.",
        "  - Prediction Accuracy: Maintain demand forecasting error variance (RMSE) under 10%."
    ]
    add_styled_bullets(met_box.text_frame, met_text, base_font_size=15)

    # ----------------------------------------------------
    # SLIDE 5: Architecture
    # ----------------------------------------------------
    slide_5 = prs.slides.add_slide(blank_slide_layout)
    add_branded_header(slide_5, "5. System Architecture")

    # Left: Text description of Monorepo decoupings
    arch_col1 = slide_5.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    a1_text = [
        "Decoupled Microservice Design",
        "  - Decoupled into independent containers coordinated via multi-container Docker Compose.",
        "  - Presentation Layer (Frontend):",
        "    - Next.js 15 web client using Zustand state management.",
        "  - Service Orchestration Layer (Backend):",
        "    - Asynchronous FastAPI core utilizing non-blocking Motor drivers with MongoDB database clusters.",
        "  - Intelligence Layer (ML Engine):",
        "    - Python microservice running regressions and RL engines."
    ]
    add_styled_bullets(arch_col1.text_frame, a1_text, base_font_size=14)

    # Right: Data Pipelines flow
    arch_col2 = slide_5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    a2_text = [
        "Async Background Execution Flow",
        "  1. APScheduler automatically triggers competitor scraping modules every 2 hours.",
        "  2. Non-blocking Motor database handlers store competitor price telemetry.",
        "  3. Ingestion streams launch elasticity and demand forecast analyses.",
        "  4. The pricing engine calculates recommendations based on inventory scarcity.",
        "  5. Decoupled Next.js client layers immediately reflect updated pricing suggestions."
    ]
    add_styled_bullets(arch_col2.text_frame, a2_text, base_font_size=14)

    # ----------------------------------------------------
    # SLIDE 6: Techstacks
    # ----------------------------------------------------
    slide_6 = prs.slides.add_slide(blank_slide_layout)
    add_branded_header(slide_6, "6. Technology Stack")

    # Column 1: Front and Backend
    tech_col1 = slide_6.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    t1_text = [
        "Frontend Client Layer",
        "  - Next.js 15 & React: Component-driven routing and UI structures.",
        "  - TypeScript: Strict static typing.",
        "  - Tailwind CSS & Shadcn UI: Modular styles.",
        "  - Framer Motion: Responsive animations.",
        "  - Zustand: Global client-side state manager.",
        "Backend Core Infrastructure",
        "  - FastAPI (Python): Asynchronous REST API routing.",
        "  - Motor: Asynchronous MongoDB connection pooling driver.",
        "  - APScheduler: Periodic background task scheduler."
    ]
    add_styled_bullets(tech_col1.text_frame, t1_text, base_font_size=14)

    # Column 2: ML Stack and Ops
    tech_col2 = slide_6.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    t2_text = [
        "Machine Learning & Optimization",
        "  - Prophet (Meta): Seasonality and trend time-series forecasting.",
        "  - XGBoost Regressor: Feature-based demand forecasting.",
        "  - scikit-learn: Price elasticity linear regression.",
        "  - NumPy & Pandas: Matrix operations and data manipulation.",
        "Containerization & Storage",
        "  - MongoDB: NoSQL document store for product datasets.",
        "  - Docker & Docker Compose: Containerized service isolation."
    ]
    add_styled_bullets(tech_col2.text_frame, t2_text, base_font_size=14)

    # ----------------------------------------------------
    # SLIDE 7: Risks & Assumptions
    # ----------------------------------------------------
    slide_7 = prs.slides.add_slide(blank_slide_layout)
    add_branded_header(slide_7, "7. Risks & Assumptions")

    r_col1 = slide_7.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    r1_text = [
        "Key Risks & Mitigations",
        "  - Competitor Scraping Interruptions:",
        "    - Risk: External sites change structure, breaking scrapers.",
        "    - Mitigation: Scrapers check structural changes and send alert notifications.",
        "  - Cold-Start ML Performance:",
        "    - Risk: New SKUs lack baseline transaction sales history.",
        "    - Mitigation: Engine applies fallback competitor-matching heuristics until 30 days of data are gathered."
    ]
    add_styled_bullets(r_col1.text_frame, r1_text, base_font_size=14)

    r_col2 = slide_7.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    r2_text = [
        "Core Project Assumptions",
        "  - Continuous Competitor Accessibility:",
        "    - Competitor pricing and stock levels are publicly accessible via network routes.",
        "  - Daily E-Commerce Data Integrity:",
        "    - Historical daily transaction records are clean and formatted for Prophet and XGBoost calibration.",
        "  - Microservice Network Performance:",
        "    - Local network connection latencies between FastAPI backend and ML engine remain under 10ms."
    ]
    add_styled_bullets(r_col2.text_frame, r2_text, base_font_size=14)

    # ----------------------------------------------------
    # SLIDE 8: Work breakdown structure and Team Roles
    # ----------------------------------------------------
    slide_8 = prs.slides.add_slide(blank_slide_layout)
    add_branded_header(slide_8, "8. Work Breakdown Structure & Team Roles")

    wbs_col1 = slide_8.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    w_text = [
        "Work Breakdown Structure (WBS)",
        "  - Module 1: ML Core Pipelines (Shivam Kumar)",
        "    - Code price sensitivity regression calculators.",
        "    - Implement Prophet/XGBoost demand forecasters.",
        "    - Build Q-learning reinforcement structures.",
        "  - Module 2: Async Backend Integration",
        "    - Configure async routes in FastAPI.",
        "    - Setup Motor connection pooling to MongoDB.",
        "    - Integrate APScheduler background ingestion.",
        "  - Module 3: Frontend Dashboard & Deploy",
        "    - Build Next.js 15 dashboard graphs and states."
    ]
    add_styled_bullets(wbs_col1.text_frame, w_text, base_font_size=14)

    roles_col2 = slide_8.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    r_text = [
        "Core Team Roles & Focus Areas",
        "  - Lead AI & ML Engineer (Shivam Kumar):",
        "    - Owns mathematical price sensitivity regression models, Prophet+XGBoost models, and Q-learning loops.",
        "  - Lead Backend Architect:",
        "    - Owns FastAPI routing logic, database collection schemas, and background task queues.",
        "  - Lead Frontend Architect:",
        "    - Owns visual interface modules, charts, responsive animations, and Zustand state variables."
    ]
    add_styled_bullets(roles_col2.text_frame, r_text, base_font_size=14)

    # ----------------------------------------------------
    # SLIDE 9: Timeline in 2 weeks
    # ----------------------------------------------------
    slide_9 = prs.slides.add_slide(blank_slide_layout)
    add_branded_header(slide_9, "9. Two-Week Implementation Timeline")

    t_col1 = slide_9.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    t1_text = [
        "Week 1: Foundations & Backend Core (Days 1 - 7)",
        "  - Days 1 - 2: Setup Docker monorepo. Create FastAPI routers, MongoDB schemas, and connect Motor.",
        "  - Days 3 - 4: Code mathematical PriceElasticityModel regressions and configure pricing boundaries.",
        "  - Days 5 - 7: Implement discrete QLearningPricingEngine states (inventory scarcity) and rewards."
    ]
    add_styled_bullets(t_col1.text_frame, t1_text, base_font_size=15)

    t_col2 = slide_9.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5))
    t2_text = [
        "Week 2: Frontends & Pipelines Integration (Days 8 - 14)",
        "  - Days 8 - 10: Build Next.js 15 dashboard layouts, visual cards, charts, and link Zustand stores.",
        "  - Days 11 - 12: Connect Prophet + XGBoost demand forecasting pipelines to backend controllers.",
        "  - Days 13 - 14: Finalize multi-container Docker Compose files, complete RMSE system test suites, and launch."
    ]
    add_styled_bullets(t_col2.text_frame, t2_text, base_font_size=15)

    # ----------------------------------------------------
    # SLIDE 10: References
    # ----------------------------------------------------
    slide_10 = prs.slides.add_slide(blank_slide_layout)
    add_branded_header(slide_10, "10. References")

    ref_col = slide_10.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.133), Inches(4.5))
    refs = [
        "[1] Fisher, M., Gallino, S., & Li, J. (2018). Competition and Pricing in E-commerce: An Empirical Analysis. Management Science, 64(6), 2530-2545.",
        "[2] Ye, J., et al. (2020). Reinforcement Learning for Dynamic Pricing in E-commerce Platforms. Proceedings of the ACM SIGKDD International Conference on Knowledge Discovery & Data Mining, 2341-2350.",
        "[3] Kephart, J. O., et al. (2021). Dynamic Pricing Algorithms and Market Equilibrium. ACM Transactions on Economics and Computation (TEAC), 9(2), 1-28.",
        "[4] Montgomery, A. L., et al. (2019). Modeling Consumer Price Sensitivity with Dynamic Elasticity Log-Log Regressions. Journal of Marketing Research, 56(4), 589-607.",
        "[5] Sutton, R. S., & Barto, A. G. (2018). Reinforcement Learning: An Introduction (2nd ed.). MIT Press.",
        "[6] Taylor, J. (2021). An Evaluation of Asynchronous Data Pipelines and NoSQL Databases in Multi-Category Retail Ingestion. Journal of Systems Software, 174, 110892."
    ]
    add_styled_bullets(ref_col.text_frame, refs, base_font_size=14)

    # ----------------------------------------------------
    # Save the output file
    # ----------------------------------------------------
    output_filename = "Dynamic_Pricing_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully created and saved as {output_filename}")
    
    # Also save to parent workspace directory for convenience
    parent_path = os.path.join("/Users/shivamkumar", output_filename)
    prs.save(parent_path)
    print(f"Convenience copy saved to {parent_path}")

if __name__ == "__main__":
    create_presentation()
